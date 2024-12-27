from flask import Flask, request, jsonify, send_file, render_template
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

app = Flask(__name__)

@app.route('/')
def home():
    return render_template('index.html')  # Pastikan file index.html ada di folder templates

@app.route('/generate-ppt', methods=['POST'])
def generate_ppt():
    data = request.json  # Menerima data dari frontend
    
    # Validasi input
    if not data or 'slides' not in data or not isinstance(data['slides'], list) or len(data['slides']) == 0:
        return jsonify({"error": "No slides data provided"}), 400
    
    slides = data.get('slides', [])
    
    # Membuat file PowerPoint
    ppt = Presentation()
    
    # Fungsi untuk menambahkan slide dengan desain lebih menarik
    def add_slide_with_design(ppt, title, content):
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Menggunakan layout kosong (tidak ada judul)

        # Menambahkan latar belakang dengan warna gelap
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(44, 62, 80)  # Latar belakang gelap biru kehijauan

        # Menambahkan kotak teks untuk subjudul
        left = Inches(1)
        top = Inches(0.5)
        width = Inches(8.5)
        height = Inches(1)

        textbox_title = slide.shapes.add_textbox(left, top, width, height)
        text_frame_title = textbox_title.text_frame
        text_frame_title.text = title
        
        # Menambahkan desain teks untuk subjudul
        for paragraph in text_frame_title.paragraphs:
            paragraph.font.size = Pt(32)  # Ukuran font yang lebih besar untuk subjudul
            paragraph.font.color.rgb = RGBColor(236, 240, 241)  # Warna font putih
            paragraph.font.bold = True  # Membuat teks lebih tebal
            paragraph.font.name = 'Arial'

        # Menambahkan kotak teks untuk konten
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8.5)
        height = Inches(6.5)

        textbox_content = slide.shapes.add_textbox(left, top, width, height)
        text_frame_content = textbox_content.text_frame
        text_frame_content.text = content
        
        # Menambahkan desain teks untuk konten
        for paragraph in text_frame_content.paragraphs:
            paragraph.font.size = Pt(28)  # Ukuran font yang lebih besar
            paragraph.font.color.rgb = RGBColor(236, 240, 241)  # Warna font putih
            paragraph.font.bold = False  # Konten tidak tebal
            paragraph.font.name = 'Arial'

        # Menambahkan efek gradasi pada latar belakang
        fill = slide.background.fill
        fill.gradient()
        stops = fill.gradient_stops
        stop1 = stops[0]
        stop1.position = 0.0
        stop1.color.rgb = RGBColor(34, 49, 63)  # Warna gelap pertama

        stop2 = stops[1]
        stop2.position = 1.0
        stop2.color.rgb = RGBColor(44, 62, 80)  # Warna gelap kedua

        return slide

    # Menambahkan slide dengan desain menarik ke PPT
    for slide_data in slides:
        title = slide_data.get('title', 'Judul Slide')  # Gunakan 'Judul Slide' sebagai default jika tidak ada
        content = slide_data.get('content', '')  # Gunakan string kosong sebagai default jika tidak ada konten
        add_slide_with_design(ppt, title, content)

    # Menyimpan PPT ke file
    output_path = "generated_ppt.pptx"
    ppt.save(output_path)
    
    return send_file(output_path, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
